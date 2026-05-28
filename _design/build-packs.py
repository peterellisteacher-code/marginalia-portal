"""
build-packs.py -- assemble Stage 1 Issues Study cache packs for the chamber.

Reads source PDF/DOCX/PPTX files from website/readings_source/ and website/text_packs/,
extracts each to an inspectable .md intermediate under
website/readings_source/extracted/<reading-id>.md, then concatenates
the relevant readings into website/data/packs/<pack_id>.txt with
=== Author -- Title === separators -- the format chat.js expects.

Run from the repo root (website/ as cwd parent), or from anywhere with the
WEBSITE_ROOT constant pointing to the right place.

Re-runnable: overwrites both .md intermediates and pack .txt files every time.

Dependencies: pdfplumber, python-docx, python-pptx (all already installed).
"""

from __future__ import annotations

import re
import sys
from pathlib import Path

import pdfplumber
import fitz  # PyMuPDF
import docx
from pptx import Presentation


WEBSITE_ROOT = Path(__file__).resolve().parent.parent  # .../4 - Issues Study/website/


# Per-pack scope text and reading list.
# Tuple format: (reading_id, author, title, path-relative-to-website-root)
PACK_MAP: dict[str, dict] = {
    "stage1_existentialism": {
        "scope": (
            "Existentialist tradition (Sartre, de Beauvoir, Camus, Kierkegaard). "
            "Use these passages when the student is working on authenticity, "
            "radical freedom, bad faith, the absurd, situated freedom, or the "
            "ethical leap. Cite specific passages by reading title; do not "
            "summarise unprompted."
        ),
        "readings": [
            ("sartre", "Jean-Paul Sartre", "Sartre on Existentialism",
             "readings_source/Sartre.pdf"),
            ("sartre-republic-silence", "Jean-Paul Sartre", "The Republic of Silence",
             "readings_source/Existentialism/Sartre/Sartre - Republic of Silence.pdf"),
            ("beauvoir", "Simone de Beauvoir", "On Beauvoir's Ethics of Ambiguity",
             "readings_source/Simone de Beauvoir.pdf"),
            ("camus", "Albert Camus", "On the Absurd",
             "readings_source/Albert Camus.pdf"),
            ("existentialism-four-thinkers", "Overview", "Existentialism -- Four Thinkers",
             "readings_source/Existentialism Four Thinkers.pdf"),
            # Kierkegaard enrichment dropped 2026-05-28: text_packs source PDF
            # yielded only ~200 words across 145 pages (encoding issue).
            # Existentialism pack is already substantive at ~14K tokens.
        ],
    },
    "stage1_virtue_compassion": {
        "scope": (
            "Virtue ethics centred on Murdoch's account of moral attention and "
            "Nussbaum's account of compassion as cognition. Gilligan adds a "
            "feminist care-ethics counterweight. Use when the student is working "
            "on what it is to be a good person, the role of emotion in moral "
            "judgement, or the limits of consequentialist / Kantian accounts."
        ),
        "readings": [
            ("murdoch", "Iris Murdoch", "The Sovereignty of Good",
             "readings_source/Murdoch - Sovereignty of Good.pdf"),
            ("nussbaum", "Martha Nussbaum", "Compassion and Reason",
             "readings_source/Nussbaum - Compassion and Reason.pdf"),
            ("gilligan", "Carol Gilligan",
             "In a Different Voice (enrichment -- care ethics)",
             "text_packs/1982_in_a_different_voice_carol_gilligan.text.pdf"),
        ],
    },
    "stage1_religion_ethics": {
        "scope": (
            "Religion, moral binding, and moral foundations. Haidt's argument: "
            "religion's primary function is moral binding rather than metaphysical "
            "claim. The VCE chapter gives explicit Euthyphro-dilemma framing as "
            "a non-Haidt counterweight. Use for questions about religion and "
            "morality, divine command theory, secular ethics, and moral "
            "foundations theory."
        ),
        "readings": [
            ("haidt-bare", "Jonathan Haidt",
             "On Religion -- Bare Minimum (short version)",
             "readings_source/Haidt on Religion BARE MINIMUM.pdf"),
            ("haidt-religion", "Jonathan Haidt", "On Religion -- Standard Version",
             "readings_source/Haidt on Religion.pdf"),
            ("haidt-more", "Jonathan Haidt", "On Religion -- Extended Version",
             "readings_source/Haidt on Religion MORE THAT MINIMUM.pdf"),
            ("vce-god", "VCE Philosophy textbook",
             "On the Existence and Nature of God (enrichment -- Euthyphro framing)",
             "text_packs/chapter_vce_on_the_existence_and_nature_of_god.text.pdf"),
        ],
    },
    "stage1_aesthetics": {
        "scope": (
            "Aesthetics -- what is an artwork and where does its meaning live? "
            "Wimsatt and Beardsley's intentional fallacy is the starting move; "
            "Freeland adds the question of whether art constitutes a form of "
            "moral knowledge. Use for questions about interpretation, "
            "intention, and the ontology of artworks."
        ),
        "readings": [
            ("intentional-fallacy", "W. K. Wimsatt and Monroe Beardsley",
             "The Intentional Fallacy",
             "readings_source/Art - Wimsatt and Beardsley - the Intentional Fallacy.pdf"),
            ("intentional-unintentional", "Companion essay",
             "The Intentional and the Unintentional",
             "readings_source/The Intentional and the Unintentional.pdf"),
            ("freeland", "Cynthia Freeland",
             "Art and Moral Knowledge (enrichment)",
             "text_packs/freeland_1997_art_and_moral_knowledge.text.pdf"),
        ],
    },
    "stage1_mind_simulation": {
        "scope": (
            "Philosophy of mind, personal identity, the simulation argument, and "
            "death as the ending of mind. Identity theory and functionalism are "
            "the canonical materialist accounts; dualism is the contrast. Simon "
            "Beck adds personal-identity-over-time thought experiments. "
            "'When We Die' belongs here because it treats death as a question "
            "about personal identity. Use for mind-body questions, what makes "
            "us the same person over time, simulation worries, and the "
            "philosophy of death."
        ),
        "readings": [
            ("identity-functionalism", "Overview",
             "Identity Theory and Functionalism",
             "readings_source/Identity Theory and Functionalism.pdf"),
            ("problems-identity", "Overview", "Problems for Identity Theory",
             "readings_source/Problems for Identity Theory.pdf"),
            ("dualism", "Overview", "Dualism",
             "readings_source/Dualism.pdf"),
            ("simulation", "Class slides", "Simulation Theory",
             "readings_source/Simulation Theory.pptx"),
            ("what-is-human", "Class slides", "What is a Human Being?",
             "readings_source/What is a Human Being.pptx"),
            ("when-we-die", "Overview", "When We Die",
             "readings_source/When we die.pdf"),
            # Simon Beck enrichment dropped 2026-05-28: text_packs source PDF
            # is a 1KB stub with only ~200 words of usable text. Pack is
            # already ~15K tokens without it.
        ],
    },
    "lab_applied_normative_ethics": {
        "scope": (
            "Applied normative ethics -- personhood, abortion, and the ethics "
            "of human enhancement. Marquis, Thomson, and Singer form the "
            "abortion-ethics triangle (Future-Like-Ours, bodily autonomy, "
            "preference utilitarianism). Brave New World and the Genetic "
            "Supermarket open the enhancement question. Use for questions "
            "about moral status, what we owe to potential persons, and the "
            "ethics of designing future people."
        ),
        "readings": [
            ("singer-abortion", "Peter Singer", "On Abortion",
             "readings_source/Singer Abortion.pdf"),
            ("marquis", "Don Marquis",
             "Why Abortion is Immoral (Future-Like-Ours)",
             "readings_source/marquis.pdf"),
            ("thomson", "Judith Jarvis Thomson", "A Defense of Abortion",
             "readings_source/judith.pdf"),
            ("brave-new-world-ch16", "Aldous Huxley",
             "Brave New World, Chapter 16 -- Being Human",
             "readings_source/Brave New World Chapter 16 - Being Human.pdf"),
            ("genetic-supermarket", "Supplementary",
             "The Genetic Supermarket",
             "readings_source/Genetic Supermarket.pdf"),
        ],
    },
}


# ---------- extraction ----------

def extract_pdf(path: Path) -> str:
    """Primary extractor: PyMuPDF (fitz) -- handles fonts, multi-column,
    and ligatures better than pdfplumber for the prose-heavy academic PDFs
    in this corpus. Falls back to pdfplumber only if PyMuPDF returns empty
    (signals a scanned/image PDF that needs OCR override anyway).
    """
    doc = fitz.open(path)
    try:
        pages = [page.get_text() for page in doc]
    finally:
        doc.close()
    text = "\n\n".join(pages).strip()
    if text:
        return text
    # Fallback to pdfplumber (rare; usually empty here too means image-only)
    with pdfplumber.open(path) as pdf:
        return "\n\n".join((p.extract_text() or "") for p in pdf.pages)


def extract_docx(path: Path) -> str:
    doc = docx.Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_pptx(path: Path) -> str:
    """Slide text only -- skip speaker notes per the build spec."""
    prs = Presentation(path)
    chunks: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        lines: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    lines.append(t)
        if lines:
            chunks.append(f"[Slide {i}]\n" + "\n".join(lines))
    return "\n\n".join(chunks)


def clean(text: str) -> str:
    """Light cleanup of common PDF/DOCX extraction artifacts."""
    if not text:
        return ""
    # Fix end-of-line hyphenation: "exis-\ntence" -> "existence"
    text = re.sub(r"(\w)-\n(\w)", r"\1\2", text)
    # Collapse runs of internal whitespace
    text = re.sub(r"[ \t]{2,}", " ", text)
    # Collapse 3+ blank lines to 2
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Strip whitespace around lines
    text = "\n".join(line.rstrip() for line in text.splitlines())
    return text.strip()


def extract_any(path: Path) -> str:
    ext = path.suffix.lower()
    if path.name.lower().endswith(".text.pdf") or ext == ".pdf":
        return clean(extract_pdf(path))
    if ext == ".docx":
        return clean(extract_docx(path))
    if ext == ".pptx":
        return clean(extract_pptx(path))
    raise ValueError(f"unsupported extension: {path}")


def load_override(reading_id: str) -> str | None:
    """If a hand-OCR'd / manually-edited .md exists under
    _design/ocr-overrides/<reading_id>.md, use its body as the canonical text.
    Strips the first H1 heading line if present (we re-emit our own header
    in the pack file).
    """
    override = WEBSITE_ROOT / "_design" / "ocr-overrides" / f"{reading_id}.md"
    if not override.exists():
        return None
    raw = override.read_text(encoding="utf-8")
    # Drop the leading "# Title" line if present, plus any blank line that follows
    lines = raw.splitlines()
    if lines and lines[0].startswith("# "):
        lines = lines[1:]
        while lines and not lines[0].strip():
            lines = lines[1:]
    return "\n".join(lines).strip()


# ---------- build ----------

def main() -> int:
    extracted_dir = WEBSITE_ROOT / "readings_source" / "extracted"
    packs_dir = WEBSITE_ROOT / "data" / "packs"
    extracted_dir.mkdir(parents=True, exist_ok=True)
    packs_dir.mkdir(parents=True, exist_ok=True)

    report: list[tuple[str, str, int, str]] = []  # (pack, reading, words, status)
    missing: list[str] = []

    for pack_id, spec in PACK_MAP.items():
        pack_parts: list[str] = [
            f"=== PACK: {pack_id} ===\n\n{spec['scope']}",
        ]
        for reading_id, author, title, rel_path in spec["readings"]:
            src = WEBSITE_ROOT / rel_path
            if not src.exists():
                missing.append(str(src))
                report.append((pack_id, reading_id, 0, "MISSING"))
                continue
            override_body = load_override(reading_id)
            if override_body is not None:
                body = override_body
                source_note = "ocr-override"
            else:
                try:
                    body = extract_any(src)
                except Exception as exc:
                    print(f"!! extract failed for {src}: {exc}", file=sys.stderr)
                    report.append((pack_id, reading_id, 0, f"FAILED: {exc}"))
                    continue
                source_note = "auto-extract"

            words = len(body.split())

            # Per-reading .md intermediate (inspectable, editable, re-runnable)
            md_path = extracted_dir / f"{reading_id}.md"
            md_path.write_text(
                f"# {title}\n\n_{author}_\n\nSource: `{rel_path}` "
                f"(via {source_note})\n\n---\n\n{body}\n",
                encoding="utf-8",
            )

            # Append to pack body
            pack_parts.append(f"=== {author} -- {title} ===\n\n{body}")
            report.append((pack_id, reading_id, words, source_note))

        pack_text = "\n\n".join(pack_parts) + "\n"
        (packs_dir / f"{pack_id}.txt").write_text(pack_text, encoding="utf-8")

    # ----- report -----
    print("=" * 64)
    print("BUILD REPORT")
    print("=" * 64)
    by_pack: dict[str, list[tuple[str, int, str]]] = {}
    for pack_id, reading_id, words, status in report:
        by_pack.setdefault(pack_id, []).append((reading_id, words, status))

    grand_total_words = 0
    OK_STATUSES = {"auto-extract", "ocr-override"}
    for pack_id, items in by_pack.items():
        total_words = sum(w for _, w, st in items if st in OK_STATUSES)
        # Rough English-prose estimate: 0.75 words per token
        est_tokens = int(total_words / 0.75)
        cache_status = (
            "CACHES" if est_tokens >= 4096
            else "TOO SMALL TO CACHE (<4096 tokens)"
        )
        grand_total_words += total_words
        print(f"\n[{pack_id}] {total_words:,} words, ~{est_tokens:,} tokens -- {cache_status}")
        for reading_id, words, status in items:
            marker = "  " if status in OK_STATUSES else "!!"
            print(f"  {marker} {reading_id:<30} {words:>7,} words  ({status})")

    print(f"\n{'=' * 64}")
    print(f"TOTAL: {grand_total_words:,} words across all packs "
          f"(~{int(grand_total_words / 0.75):,} tokens)")
    print(f"Packs written to: {packs_dir}")
    print(f"Intermediates written to: {extracted_dir}")
    if missing:
        print(f"\nMISSING source files ({len(missing)}):")
        for m in missing:
            print(f"  - {m}")
    print("=" * 64)
    return 1 if missing else 0


if __name__ == "__main__":
    sys.exit(main())
